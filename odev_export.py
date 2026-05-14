"""Export Excel and PowerPoint deliverables."""
from __future__ import annotations

from pathlib import Path
from typing import Dict, Optional

import pandas as pd


def export_excel(
    path: str,
    raw_adj: pd.DataFrame,
    events: pd.DataFrame,
    corr: Optional[pd.DataFrame] = None,
    corr_by_window: Optional[Dict[str, pd.DataFrame]] = None,
) -> None:
    p = Path(path)
    with pd.ExcelWriter(p, engine="openpyxl") as w:
        raw_adj.to_excel(w, sheet_name="adj_close_wide")
        events.to_excel(w, sheet_name="events_features", index=False)
        if corr is not None:
            corr.to_excel(w, sheet_name="corr_matrix")
        if corr_by_window:
            for key, mat in corr_by_window.items():
                if mat is None or mat.empty:
                    continue
                sheet = f"corr_{key}"[:31]
                mat.to_excel(w, sheet_name=sheet)


def export_presentation(
    path: str,
    title: str,
    bullets: list[tuple[str, list[str]]],
) -> None:
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = "Earnings sonrası getiri tahmini — özet"

    bullet_layout = prs.slide_layouts[1]
    for heading, lines in bullets:
        s = prs.slides.add_slide(bullet_layout)
        s.shapes.title.text = heading
        body = s.shapes.placeholders[1].text_frame
        body.clear()
        for i, line in enumerate(lines):
            if i == 0:
                p = body.paragraphs[0]
            else:
                p = body.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = Pt(18)
    prs.save(path)
